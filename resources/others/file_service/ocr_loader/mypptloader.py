from typing import List

import tqdm
from langchain_core.documents import Document
from langchain_community.document_loaders.unstructured import UnstructuredFileLoader


class RapidOCRPPTLoader(UnstructuredFileLoader):
    
    """
    重写 _get_elements 方法，返回 Document 对象列表
    """
    def _get_elements(self) -> List:

        """
        自定义PPT文档加载器（继承LangChain通用加载器）
        核心功能：
        1. 提取PPT幻灯片中的文本框文字、表格文字
        2. 提取PPT中的图片，调用RapidOCR识别图片文字
        3. 处理组合形状（ShapeType=6）：递归提取子形状内容
        4. 按幻灯片内形状的位置（从上到下、从左到右）排序提取，保证文字顺序符合视觉逻辑
        5. 拼接所有文字，返回LangChain兼容的文档格式
        """
        def ppt2text(filepath):
            from io import BytesIO
            import numpy as np
            from PIL import Image
            from pptx import Presentation
            from rapidocr_onnxruntime import RapidOCR

            ocr = RapidOCR()
            prs = Presentation(filepath)
            resp = ""

            """
            递归提取PPT幻灯片中的文本、表格、图片等元素
            """
            def extract_text(shape):
                nonlocal resp
                # 处理文本框形状：提取文本框内的文字
                if shape.has_text_frame:
                    resp += shape.text.strip() + "\n"
                # 处理表格形状：提取表格内的文字
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                resp += paragraph.text.strip() + "\n"
                # 处理图片形状：调用RapidOCR识别图片文字
                if shape.shape_type == 13:  # 13 表示图片
                    image = Image.open(BytesIO(shape.image.blob))
                    result, _ = ocr(np.array(image))
                    if result:
                        ocr_result = [line[1] for line in result]
                        resp += "\n".join(ocr_result)
                # 处理组合形状：ShapeType=6 是python-pptx中组合形状的固定类型值，递归提取子形状内容
                elif shape.shape_type == 6:  # 6 表示组合
                    for child_shape in shape.shapes:
                        extract_text(child_shape)

            # 初始化进度条
            b_unit = tqdm.tqdm(
                total=len(prs.slides), desc="RapidOCRPPTLoader slide index: 1"
            )
            # 遍历所有幻灯片
            for slide_number, slide in enumerate(prs.slides, start=1):
                b_unit.set_description(
                    "RapidOCRPPTLoader slide index: {}".format(slide_number)
                )
                b_unit.refresh()
                sorted_shapes = sorted(
                    slide.shapes, key=lambda x: (x.top, x.left)
                )  # 从上到下、从左到右遍历
                for shape in sorted_shapes:
                    extract_text(shape)
                b_unit.update(1)
            return resp

        text = ppt2text(self.file_path)
        # return [Document(
        #     page_content=text.strip(),  # 提取的纯文字（去除首尾空格）
        #     metadata={
        #         "source": self.file_path,  # 文件路径（溯源用）
        #         "file_type": "pptx",       # 标记文件类型（便于后续处理）
        #         "processed_by": "RapidOCRPPTLoader"  # 标记处理加载器
        #     }
        # )]
        from unstructured.partition.text import partition_text
        return partition_text(text=text, **self.unstructured_kwargs)


if __name__ == "__main__":
    from pathlib import Path
    file_path = Path(__file__).parent.parent.parent / "test" / "samples" / "ocr_test.pptx"
    loader = RapidOCRPPTLoader(file_path=str(file_path))
    docs = loader.load()
    print(docs)